from pptx import Presentation
import os
import json

pptx_path = 'Solar_System_Presentation.pptx'
prs = Presentation(pptx_path)
file_size = os.path.getsize(pptx_path)
slide_count = len(prs.slides)

print("=== PRESENTATION METADATA ===")
print("File:", pptx_path)
print("File size: {} bytes ({:.1f} KB)".format(file_size, file_size/1024))
print("Total slides:", slide_count)
print()
print("=== SLIDE DETAILS ===")

slides_with_notes = 0
total_notes_chars = 0

for idx, slide in enumerate(prs.slides):
    print("Slide {}:".format(idx + 1))
    print("  Layout:", slide.slide_layout.name)
    print("  Shapes:", len(slide.shapes))
    
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text
        if notes.strip():
            slides_with_notes += 1
            total_notes_chars += len(notes)
            print("  Speaker notes: YES ({} chars)".format(len(notes)))
        else:
            print("  Speaker notes: NO")
    else:
        print("  Speaker notes: NO")
    print()

print("=== VERIFICATION CHECKLIST ===")
print("[OK] File exists: YES")
print("[{}] Contains slides: {} slides".format("OK" if slide_count > 0 else "FAIL", slide_count))
print("[{}] Meets 8+ slide requirement: {}".format("OK" if slide_count >= 8 else "FAIL", slide_count >= 8))
print("[{}] Slides with speaker notes: {}/{}".format("OK" if slides_with_notes > 0 else "WARN", slides_with_notes, slprint("[{}] Slides with speakal print("[{}] Slides w {print("[{}] Slides wal_print("[{}] Slides with speaker / max(1, slides_with_notes)))
